from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation
from PIL import Image
import os


def create_title_slide(prs, slide_title, title_font_size=None):
    """
    创建标题幻灯片，并支持动态优化排版（分行、调整字体大小）。
    遇到特殊分隔符（如“：”“: ”“——”等）自动换行。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 添加空白幻灯片

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 动态调整字体大小
    if title_font_size is None:
        title_font_size = adjust_title_font_size(len(slide_title))

    # 每行最大字符数根据字体大小动态调整
    max_chars_per_line = get_max_chars_for_font_size(title_font_size)

    # 分行逻辑：根据最大字符数和特殊分隔符分行
    lines = split_title_into_lines(slide_title, max_chars_per_line)

    # 计算整体标题的高度
    line_height = Pt(title_font_size + 10)  # 每行的高度，包括行间距
    total_height = line_height * len(lines)

    # 垂直居中计算
    title_top = (slide_height - total_height) / 2

    # 添加标题框
    title_box = slide.shapes.add_textbox(Inches(0.5), title_top, slide_width - Inches(1), total_height)
    title_frame = title_box.text_frame
    _clear_text_frame(title_frame)

    # 填充标题内容
    for line in lines:
        paragraph = title_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(title_font_size)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

    return slide

def get_max_chars_for_font_size(font_size):
    """
    根据字体大小返回每行最大字符数（中文/英文）。
    """
    max_chars_mapping = {
        48: {"zh": 11, "en": 22},
        44: {"zh": 13, "en": 26},
        40: {"zh": 15, "en": 30},
        36: {"zh": 17, "en": 34}
    }
    return max_chars_mapping.get(font_size, {"zh": 11, "en": 22})  # 默认返回 48 号字的限制

def adjust_title_font_size(title_length):
    """
    根据标题长度动态调整字体大小。最大字体大小为 48，最小为 36。
    """
    if title_length <= 11:  # 一行即可容纳
        return 48
    elif title_length <= 13:
        return 44
    elif title_length <= 15:
        return 40
    else:
        return 36

def split_title_into_lines(title, max_chars_per_line):
    """
    根据特殊分隔符和最大字符数分行标题。
    :param title: 标题文本
    :param max_chars_per_line: 每行最大字符数（中文/英文）
    :return: 分行后的标题列表
    """
    special_separators = ["：", ": ", "——"]  # 特殊分隔符
    lines = []
    current_line = ""

    while title:
        # 优先根据特殊分隔符分行
        for sep in special_separators:
            if sep in title:
                index = title.index(sep) + len(sep)
                if len(current_line) + index <= max_chars_per_line["zh"]:  # 中文字符限制
                    current_line += title[:index]
                    title = title[index:]
                else:
                    lines.append(current_line.strip())
                    current_line = ""
                break
        else:
            # 如果没有特殊分隔符，用字符数限制分行
            if len(current_line) + len(title) <= max_chars_per_line["zh"]:
                current_line += title
                title = ""
            else:
                remaining_space = max_chars_per_line["zh"] - len(current_line)
                current_line += title[:remaining_space]
                title = title[remaining_space:]

        # 如果当前行满了，添加到结果
        if len(current_line) >= max_chars_per_line["zh"]:
            lines.append(current_line.strip())
            current_line = ""

    # 添加最后一行
    if current_line:
        lines.append(current_line.strip())

    return lines

def create_full_text_slide(prs, slide_title, slide_text, font_size=18, language="zh"):
    """
    创建纯文字幻灯片，支持中英文混合分页。
    """
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 调用分页函数
    text_pages = paginate_text_mixed(slide_text, "full_text", font_size, language)

    for page_text in text_pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 添加空白幻灯片

        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), slide_width - Inches(1), Inches(1))
        title_frame = title_box.text_frame
        _clear_text_frame(title_frame)
        title_frame.text = slide_title

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER

        # 添加正文
        text_box_left = Inches(1)
        text_box_top = Inches(1.5)
        text_box_width = slide_width - Inches(2)
        text_box_height = slide_height - Inches(3)

        text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
        text_frame = text_box.text_frame
        _clear_text_frame(text_frame)

        for line in page_text.split("\n"):
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(font_size)
            paragraph.alignment = PP_ALIGN.LEFT

        text_frame.word_wrap = True  # 启用自动换行

    return slide


def create_text_above_image_slide(prs, slide_title, slide_text, image_path, image_title, font_size=18, language="zh"):
    """
    创建文本在上、图片在下的幻灯片，同时进一步降低图片高度，避免图片标题超出页面。
    """
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 定义图片的最大宽度和高度
    max_image_width = slide_width - Inches(2)  # 图片左右留白 1 英寸
    max_image_height = Inches(1.5)  # 降低图片高度为 1.5 英寸

    # 调用分页函数，预留图片空间后计算文本框高度
    text_box_height = slide_height - max_image_height - Inches(3.5)  # 文本框高度 = 总高度 - 图片高度 - 留白
    if text_box_height <= Inches(0):  # 确保文本框高度不为负
        raise ValueError("文本框高度不足，请减少图片占用空间或调整页面布局。")

    text_pages = paginate_text_mixed(slide_text, "text_above_image", font_size, language)

    for page_text in text_pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 添加空白幻灯片

        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), slide_width - Inches(1), Inches(1))
        title_frame = title_box.text_frame
        _clear_text_frame(title_frame)
        title_frame.text = slide_title

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(28)  # 标题字体大小
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER

        # 添加正文文本框
        text_box_left = Inches(1)
        text_box_top = Inches(1.5)
        text_box_width = slide_width - Inches(2)  # 左右留白各 1 英寸

        text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
        text_frame = text_box.text_frame
        _clear_text_frame(text_frame)

        for line in page_text.split("\n"):
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(font_size)  # 正文字体大小
            paragraph.alignment = PP_ALIGN.LEFT

        text_frame.word_wrap = True  # 启用自动换行

        # 添加图片
        if image_path and os.path.exists(image_path):
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height

            # 根据固定宽高限制计算图片大小，保持比例
            if max_image_width / aspect_ratio <= max_image_height:
                image_width = max_image_width
                image_height = image_width / aspect_ratio
            else:
                image_height = max_image_height
                image_width = image_height * aspect_ratio

            img_left = (slide_width - image_width) / 2  # 图片水平居中
            img_top = slide_height - image_height - Inches(0.5)  # 图片放置在页面下方，上面留白 0.5 英寸

            slide.shapes.add_picture(image_path, img_left, img_top, image_width, image_height)

            # 添加图片标题
            image_title_box = slide.shapes.add_textbox(
                img_left, img_top + image_height + Inches(0.1), image_width, Inches(0.5)
            )
            image_title_frame = image_title_box.text_frame
            _clear_text_frame(image_title_frame)
            image_title_frame.text = image_title

            image_title_paragraph = image_title_frame.paragraphs[0]
            image_title_paragraph.font.size = Pt(14)  # 图片标题字体大小
            image_title_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体
            image_title_paragraph.alignment = PP_ALIGN.CENTER

    return slide

def create_custom_slide(prs, slide_title, slide_text, image_path, image_title, font_size=18, language="zh"):
    """
    创建左右布局的幻灯片，左侧为文字，右侧为图片。
    根据分页规则，确保文字和图片合理布局。
    """
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 左右布局宽度定义
    text_box_width = slide_width * 0.6  # 左侧文字区域宽度（60%）
    image_box_width = slide_width * 0.4  # 右侧图片区域宽度（40%）
    text_box_height = slide_height - Inches(2)  # 上下各留白 1 英寸

    # 调用分页函数
    text_pages = paginate_text_mixed(slide_text, "custom", font_size, language)

    for page_text in text_pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 添加空白幻灯片

        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), slide_width - Inches(1), Inches(1))
        title_frame = title_box.text_frame
        _clear_text_frame(title_frame)
        title_frame.text = slide_title

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER

        # 添加文字框（左侧部分）
        text_box_left = Inches(0.5)  # 左侧距离 0.5 英寸
        text_box_top = Inches(1.5)  # 距离顶部 1.5 英寸

        text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
        text_frame = text_box.text_frame
        _clear_text_frame(text_frame)

        for line in page_text.split("\n"):
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(font_size)
            paragraph.alignment = PP_ALIGN.LEFT

        text_frame.word_wrap = True  # 启用自动换行

        # 添加图片框（右侧部分）
        if image_path and os.path.exists(image_path):
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height

            # 图片的最大宽高限制
            max_image_width = image_box_width - Inches(0.5)  # 图片右侧留白 0.5 英寸
            max_image_height = slide_height - Inches(2)  # 上下各留白 1 英寸

            # 根据固定宽高限制计算图片大小，保持比例
            if max_image_width / aspect_ratio <= max_image_height:
                image_width = max_image_width
                image_height = image_width / aspect_ratio
            else:
                image_height = max_image_height
                image_width = image_height * aspect_ratio

            img_left = text_box_width + Inches(0.5)  # 图片左侧紧贴文字右侧
            img_top = (slide_height - image_height) / 2  # 图片垂直居中

            slide.shapes.add_picture(image_path, img_left, img_top, image_width, image_height)

            # 添加图片标题
            image_title_box = slide.shapes.add_textbox(
                img_left, img_top + image_height + Inches(0.2), image_width, Inches(0.5)
            )
            image_title_frame = image_title_box.text_frame
            _clear_text_frame(image_title_frame)
            image_title_frame.text = image_title

            image_title_paragraph = image_title_frame.paragraphs[0]
            image_title_paragraph.font.size = Pt(14)
            image_title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            image_title_paragraph.alignment = PP_ALIGN.CENTER

    return slide

def get_max_lines(template_type, font_size):
    """
    根据模板类型和字体大小返回最大行数。
    :param template_type: 模板类型 ("full_text", "custom", "text_above_image")
    :param font_size: 字体大小
    :return: 最大行数
    """
    max_lines_mapping = {
        "full_text": {24: 13, 20: 15, 18: 18, 16: 20, 14: 24, 12: 28},
        "custom": {24: 10, 20: 12, 18: 14, 16: 16, 14: 20, 12: 24},
        "text_above_image": {24: 8, 20: 9, 18: 11, 16: 12, 14: 14, 12: 16},
    }
    return max_lines_mapping[template_type][font_size]


def paginate_text_mixed(text, template_type, font_size, language="zh"):
    """
    根据回车符和字符统计实现分页，支持中英文混排，避免英文单词被拆开。
    :param text: 输入文本
    :param template_type: 模板类型 ("full_text", "custom", "text_above_image")
    :param font_size: 字体大小 (例如 24, 20, 18, 16, 14, 12)
    :param language: 文本语言 ("zh" 表示中文, "en" 表示英文)
    :return: 分页后的文本列表
    """
    # 获取每行最大字符数和每页最大行数
    if template_type not in max_chars_mapping:
        raise ValueError(f"模板类型 '{template_type}' 不存在于 max_chars_mapping 中。")

    if font_size not in max_chars_mapping[template_type]:
        raise ValueError(f"字体大小 '{font_size}' 未定义在模板 '{template_type}' 中。")

    max_chars = max_chars_mapping[template_type][font_size]
    max_chars_per_line = max_chars["zh"] if language == "zh" else max_chars["en"]
    max_lines_per_page = 10  # 每页最大行数，假设固定为 10 行，可根据需要调整

    # 初始化分页结果
    pages = []
    current_page = []
    current_lines = 0

    # 按段落分割文本（基于回车符）
    paragraphs = text.split("\n")

    for paragraph in paragraphs:
        words = paragraph.split(" ")  # 英文按单词分割，中英文混排也适用
        current_line = ""
        current_width = 0

        for word in words:
            # 计算单词宽度（中文字符占 2，英文字符占 1）
            word_width = sum(1 if is_chinese(char) else 0.5 for char in word) + 1  # 加空格宽度
            if current_width + word_width > max_chars_per_line:  # 超出行宽，换行
                current_page.append(current_line.strip())  # 添加当前行到页面
                current_line = word + " "  # 开始新行
                current_width = word_width
                current_lines += 1

                # 如果当前页已满，则分页
                if current_lines >= max_lines_per_page:
                    pages.append("\n".join(current_page))
                    current_page = []
                    current_lines = 0
            else:
                # 单词加入当前行
                current_line += word + " "
                current_width += word_width

        # 添加剩余的行
        if current_line.strip():
            current_page.append(current_line.strip())
            current_lines += 1

        # 如果段落结束时当前页未满，添加空行
        if current_lines < max_lines_per_page:
            current_page.append("")
            current_lines += 1

        # 如果当前页满了，分页
        if current_lines >= max_lines_per_page:
            pages.append("\n".join(current_page))
            current_page = []
            current_lines = 0

    # 添加最后一页内容
    if current_page:
        pages.append("\n".join(current_page))

    return pages


def is_chinese(char):
    """
    判断一个字符是否是中文。
    """
    return '\u4e00' <= char <= '\u9fff'

# 全局定义 max_chars_mapping
max_chars_mapping = {
    "full_text": {
        24: {"zh": 23, "en": 46},
        20: {"zh": 28, "en": 56},
        18: {"zh": 31, "en": 62},
        16: {"zh": 35, "en": 70},
        14: {"zh": 40, "en": 80},
        12: {"zh": 46, "en": 92},
    },
    "custom": {
        24: {"zh": 14, "en": 28},
        20: {"zh": 17, "en": 34},
        18: {"zh": 19, "en": 38},
        16: {"zh": 21, "en": 42},
        14: {"zh": 24, "en": 48},
        12: {"zh": 28, "en": 56},
    },
    "text_above_image": {
        24: {"zh": 23, "en": 46},
        20: {"zh": 28, "en": 56},
        18: {"zh": 31, "en": 62},
        16: {"zh": 35, "en": 70},
        14: {"zh": 40, "en": 80},
        12: {"zh": 46, "en": 92},
    },
}

def _clear_text_frame(text_frame):
    """
    清除默认段落。
    """
    for paragraph in text_frame.paragraphs:
        text_frame._element.remove(paragraph._element)