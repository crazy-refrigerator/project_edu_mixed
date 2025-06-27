from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def create_full_text_slide(output_path, slide_title, slide_text):
    prs = Presentation(output_path) if os.path.exists(output_path) else Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 添加空白幻灯片
    slide_width = Inches(10)

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), slide_width, Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.text = slide_title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.alignment = PP_ALIGN.LEFT

    # 添加正文
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), slide_width, Inches(5))
    text_frame = text_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.add_paragraph()
    paragraph.text = slide_text
    paragraph.font.size = Pt(14)

    prs.save(output_path)

def add_title(slide, title_text, margin_left, margin_top, slide_width):
    """
    添加统一样式的标题到幻灯片
    :param slide: 当前幻灯片对象
    :param title_text: 标题文本
    :param margin_left: 标题左边距
    :param margin_top: 标题上边距
    :param slide_width: 幻灯片可用宽度
    """
    title_box = slide.shapes.add_textbox(margin_left, margin_top, slide_width, Pt(40))
    title_frame = title_box.text_frame
    title_frame.clear()  # 清除默认段落
    title_frame.text = title_text.strip()
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(28)  # 标题字体大小
    title_paragraph.font.bold = True  # 标题加粗
    title_paragraph.alignment = PP_ALIGN.LEFT  # 标题左对齐
    # title_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体