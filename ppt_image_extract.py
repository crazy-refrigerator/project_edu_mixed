import os
import re
from pptx import Presentation


def extract_images_from_pptx(pptx_path, output_dir, target_slides=None):
    """
    从 PPT 文件中提取图片并保存到指定目录，同时提取包含图号的文本。
    :param pptx_path: PPT 文件路径
    :param output_dir: 图片保存目录
    :param target_slides: 需要处理的幻灯片索引列表，None 表示处理所有幻灯片
    :return: 包含图片路径和图号的列表
    """
    # 获取 PPT 文件名（无扩展名）
    ppt_name = os.path.splitext(os.path.basename(pptx_path))[0]
    ppt_output_dir = os.path.join(output_dir, ppt_name)
    os.makedirs(ppt_output_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    results = []

    image_counter = 1
    for slide_index, slide in enumerate(prs.slides, start=1):
        # 如果有指定的幻灯片范围且当前幻灯片不在范围内，跳过
        if target_slides and slide_index not in target_slides:
            continue

        # 提取幻灯片中的所有文本
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + " "

        for shape in slide.shapes:
            if shape.shape_type == 13:  # 图片类型
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext  # 图片扩展名

                # 保存图片
                image_filename = f"slide_{slide_index}_img_{image_counter}.{image_ext}"
                image_output_path = os.path.join(ppt_output_dir, image_filename)
                with open(image_output_path, "wb") as f:
                    f.write(image_bytes)

                # 提取图号（如果存在）
                figure_number = extract_figure_number(all_text)

                # 如果没有图号，自动生成
                if not figure_number:
                    figure_number = f"Figure {image_counter}"

                # 构建相对路径
                relative_path = os.path.relpath(image_output_path, "static")

                results.append({
                    "slide": slide_index,
                    "image_index": image_counter,
                    "image_path": relative_path,
                    "figure_number": figure_number,
                })
                image_counter += 1

    return results


def extract_figure_number(text):
    """
    使用正则表达式从文本中提取图号。
    :param text: 提取的页面文本
    :return: 提取到的图号或 None
    """
    patterns = [
        r"Fig(?:ure)?\.?\s?\d+(\.\d+)?",  # 匹配 'Fig. 1.1' 或 'Figure 1.2'
        r"图\d+(\.\d+)?",                # 匹配 '图1.1' 或 '图2'
    ]
    for pattern in patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            return match.group(0)
    return None